from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
from groq import Groq
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from typing import List
from typing_extensions import Annotated
import os
import json
import random
from supabase import create_client

SUPABASE_URL = "https://pzhvqafnidrgiqklolvf.supabase.co"
SUPABASE_KEY = "sb_publishable_wBjGZ9KvjI2VnBqupse_Qw_At4oflxF"

supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

# Setup App
app = FastAPI(openapi_version="3.0.3")
print("OPENAPI VERSION:", app.openapi_version)
os.makedirs("uploads", exist_ok=True)

# Setup RAG Components
model = SentenceTransformer("all-MiniLM-L6-v2")

chroma_client = chromadb.Client(Settings(
    persist_directory="./vectorstore",
    is_persistent=True
))

collection = chroma_client.get_or_create_collection("documents")

groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# Utility Functions
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text_content = ""

    if ext == ".pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content += text + "\n"

    elif ext == ".docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            text_content += para.text + "\n"

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_content += run.text + "\n"

    else:
        raise ValueError("Format file tidak didukung")

    return text_content


def chunk_text(text, chunk_size=800):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


def ingest_file(file_path, user_id):
    text = extract_text_from_file(file_path)
    chunks = chunk_text(text)
    
    file_name = os.path.basename(file_path)
    collection.delete(where={"file_name": file_name})

    for i, chunk in enumerate(chunks):
        if chunk.strip():
            embedding = model.encode(chunk).tolist()

            collection.add(
                documents=[chunk],
                embeddings=[embedding],
                metadatas=[{
                    "file_name": file_name,
                    "chunk_index": i,
                    "user_id": user_id
                }],
                ids=[f"{file_name}_{i}"]
            )

    return text

def generate_summary_from_text(text):
    max_chars = 6000  # sedikit lebih panjang
    text_for_summary = text[:max_chars]

    prompt = f"""
Buat ringkasan komprehensif 2–3 paragraf dari dokumen berikut.
Gunakan bahasa akademik yang formal dan jelas.
Fokus pada inti pembahasan utama dan tujuan dokumen.
Langsung tuliskan saja hasil ringkasannya.

{text_for_summary}
"""

    response = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
    )

    return response.choices[0].message.content

def shuffle_options(quiz):
    for q in quiz:
        if "options" in q:
            correct_answer = q["answer"]

            correct_index = ord(correct_answer) - ord('A')
            options = q["options"]
            correct_text = options[correct_index]

            random.shuffle(options)

            new_index = options.index(correct_text)

            q["answer"] = chr(ord('A') + new_index)
            q["options"] = options

    return quiz

# Request Schema
class QuestionRequest(BaseModel):
    question: str
    file_name: str | None = None
    user_id: str
    session_id: str

class QuizRequest(BaseModel):
    file_name: str
    user_id: str
    num_questions: int = 5
    question_type: str = "mcq"

class AnswerCheckRequest(BaseModel):
    question_id: int
    user_answer: str
    quiz: list
    user_id: str
    file_name: str

class RegisterRequest(BaseModel):
    email: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str

# Register Endpoint
@app.post("/register")
def register(request: RegisterRequest):

    res = supabase.auth.sign_up({
        "email": request.email,
        "password": request.password
    })

    return {"message": "User berhasil register"}

# Login Endpoint
@app.post("/login")
def login(request: LoginRequest):

    res = supabase.auth.sign_in_with_password({
        "email": request.email,
        "password": request.password
    })

    if res.user is None:
        return {"error": "Login gagal"}

    return {
        "message": "Login berhasil",
        "user_id": res.user.id
    }

# Endpoint Upload
@app.post("/upload")
async def upload_files(
    user_id: str,
    files: Annotated[List[UploadFile], File(...)]
):

    combined_text = ""
    file_names = []

    for file in files:
        file_path = f"uploads/{file.filename}"

        with open(file_path, "wb") as buffer:
            buffer.write(await file.read())

        supabase.table("files").insert({
            "user_id": user_id,
            "file_name": file.filename
        }).execute()

        full_text = ingest_file(file_path, user_id)

        combined_text += "\n\n" + full_text
        file_names.append(file.filename)

    max_chars = 8000
    text_for_summary = combined_text[:max_chars]

    if len(files) == 1:
        prompt = f"""
Buat ringkasan komprehensif 1 paragraf dari dokumen berikut.
Gunakan bahasa akademik formal.
Langsung ke inti ringkasan.

{text_for_summary}
"""
    else:
        prompt = f"""
Berikut adalah beberapa dokumen yang diunggah.
Buat ringkasan 1 paragraf yang menjelaskan secara umum
apa tema utama dan kesamaan pembahasan dari dokumen-dokumen ini.
Langsung ke inti ringkasan.

{text_for_summary}
"""

    response = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
    )

    return {
        "message": "File berhasil diunggah",
        "uploaded_files": file_names,
        "summary": response.choices[0].message.content
    }

# Chat Endpoint
@app.post("/chat")
def chat(request: QuestionRequest):

    supabase.table("messages").insert({
        "session_id": request.session_id,
        "role": "user",
        "content": request.question
    }).execute()

    query_embedding = model.encode(request.question).tolist()

    if request.file_name:
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=6,
            where={
                "file_name": request.file_name,
                "user_id": request.user_id
            }
        )
    else:
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=6,
            where={"user_id": request.user_id}
        )

    documents = results["documents"][0]
    metadatas = results["metadatas"][0]

    context = "\n\n".join(documents)

    prompt = f"""
Jawab hanya berdasarkan konteks berikut.
Jika tidak ditemukan dalam konteks, katakan tidak ditemukan.

Konteks:
{context}

Pertanyaan:
{request.question}
"""

    #auto title
    existing_messages = supabase.table("messages") \
        .select("*") \
        .eq("session_id", request.session_id) \
        .execute()

    if len(existing_messages.data) == 1:

        title_prompt = f"""
Buat judul singkat (maks 5 kata) dari pertanyaan berikut:

"{request.question}"

Jawab hanya judulnya saja tanpa tanda kutip.
"""

        title_response = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": title_prompt}],
            model="llama-3.1-8b-instant",
        )

        title = title_response.choices[0].message.content.strip()

        supabase.table("sessions") \
            .update({"title": title}) \
            .eq("id", request.session_id) \
            .execute()

    response = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
    )

    response_text = response.choices[0].message.content

    supabase.table("messages").insert({
        "session_id": request.session_id,
        "role": "ai",
        "content": response_text
    }).execute()

    sources = []
    seen = set()

    for meta in metadatas:
        key = (meta["file_name"], meta["chunk_index"])
        if key not in seen:
            seen.add(key)
            sources.append({
                "file_name": meta["file_name"],
                "chunk_index": meta["chunk_index"]
            })

    return {
        "answer": response_text,
        "sources": sources
    }

# Endpoint Files
@app.get("/files")
def list_files(user_id: str):

    res = supabase.table("files") \
        .select("*") \
        .eq("user_id", user_id) \
        .execute()

    return res.data

    all_data = collection.get(include=["metadatas"])

    metadatas = all_data.get("metadatas", [])

    file_counts = {}

    for meta in metadatas:
        if meta and "file_name" in meta:
            file_name = meta["file_name"]
            file_counts[file_name] = file_counts.get(file_name, 0) + 1

    files = []

    for file_name, count in file_counts.items():
        files.append({
            "file_name": file_name,
            "total_chunks": count
        })

    return files

# endpoint quiz
@app.post("/generate-quiz")
def generate_quiz(request: QuizRequest):

    results = collection.query(
        query_embeddings=[model.encode("materi utama dokumen").tolist()],
        n_results=8,
        where={
            "file_name": request.file_name,
            "user_id": request.user_id
        }
    )

    documents = results["documents"][0]

    if not documents:
        return {"error": "File tidak ditemukan atau belum diupload"}

    context = "\n\n".join(documents)

    if request.question_type == "mcq":
        prompt = f"""
    Berdasarkan materi berikut, buat {request.num_questions} soal pilihan ganda.

    WAJIB:
    - Output hanya JSON
    - Jangan tambahkan penjelasan
    - Jangan pakai markdown

    Format:

    [
    {{
        "question": "...",
        "options": ["...", "...", "...", "..."],
        "answer": "A"
    }}
    ]

    Materi:
    {context}
    """
    else:
        prompt = f"""
    Berdasarkan materi berikut, buat {request.num_questions} soal essay.

    WAJIB:
    - Output hanya JSON
    - Jangan tambahkan penjelasan

    Format:

    [
    {{
        "question": "...",
        "answer": "..."
    }}
    ]

    Materi:
    {context}
    """
    response = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
    )

    raw_quiz = response.choices[0].message.content.strip()

    if raw_quiz.startswith("```"):
        raw_quiz = raw_quiz.replace("```json", "").replace("```", "").strip()

    start = raw_quiz.find("[")
    end = raw_quiz.rfind("]")

    if start != -1 and end != -1:
        raw_quiz = raw_quiz[start:end+1]

    try:
        quiz = json.loads(raw_quiz)

        for i, q in enumerate(quiz):
            q["id"] = i

    except:
        return {
            "error": "Format quiz tidak valid dari LLM",
            "raw_output": raw_quiz
        }

    if request.question_type == "mcq":
        quiz = shuffle_options(quiz)

    return {
        "quiz": quiz
    }

@app.post("/check-answer")
def check_answer(request: AnswerCheckRequest):

    question = next((q for q in request.quiz if q["id"] == request.question_id), None)

    if not question:
        return {"error": "Question tidak ditemukan"}

    correct_answer = question["answer"]

    user_embedding = model.encode(request.user_answer)
    correct_embedding = model.encode(correct_answer)

    similarity = float(
        (user_embedding @ correct_embedding) /
        ((user_embedding @ user_embedding) ** 0.5 * (correct_embedding @ correct_embedding) ** 0.5)
    )

    percentage = round(similarity * 100, 2)
    score_percent = similarity * 100

    if percentage > 85:
        feedback = "Jawaban sangat sesuai dengan konteks."
    elif percentage > 70:
        feedback = "Jawaban cukup sesuai, namun masih bisa diperbaiki."
    elif percentage > 50:
        feedback = "Jawaban kurang tepat, hanya sebagian yang relevan."
    else:
        feedback = "Jawaban tidak sesuai dengan konteks."

    supabase.table("quiz_history").insert({
        "user_id": request.user_id,
        "file_name": request.file_name,
        "question": question["question"],
        "user_answer": request.user_answer,
        "correct_answer": correct_answer,
        "similarity": similarity,
        "feedback": feedback
    }).execute()

    return {
        "similarity": percentage,
        "feedback": feedback,
        "reference_answer": correct_answer
    }

# Endpoint Delete
@app.delete("/delete-file")
def delete_file(file_name: str, user_id: str):

    # Supabase
    supabase.table("files") \
        .delete() \
        .eq("file_name", file_name) \
        .eq("user_id", user_id) \
        .execute()

    # Chroma
    collection.delete(
        where={
            "file_name": file_name,
            "user_id": user_id
        }
    )

    return {"message": "File berhasil dihapus"}

# Endpoint Session Baru
@app.post("/create-session")
def create_session(user_id: str, title: str = "New Chat"):

    res = supabase.table("sessions").insert({
        "user_id": user_id,
        "title": title
    }).execute()

    return res.data